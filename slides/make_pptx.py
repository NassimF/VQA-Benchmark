"""Generate lecturebench_slides.pptx from scratch using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── palette ──────────────────────────────────────────────────────────────────
C1     = RGBColor(0x4C, 0x72, 0xB0)   # config-1 blue
C2     = RGBColor(0xDD, 0x84, 0x52)   # config-2 orange
GREEN  = RGBColor(0x55, 0xA8, 0x68)   # eval green
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black for titles
MID    = RGBColor(0x44, 0x44, 0x66)   # body text
LIGHT  = RGBColor(0xF0, 0xF4, 0xFF)   # slide background tint
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)   # red highlight
YELLOW = RGBColor(0xFF, 0xF0, 0x80)   # table row highlight
GRAY   = RGBColor(0xCC, 0xCC, 0xCC)

# ── slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=LIGHT):
    """Fill slide background."""
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=WHITE, line=None, line_w=Pt(1)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def title_bar(slide, text, bar_h=Inches(1.0)):
    """Dark title bar across the top."""
    r = rect(slide, 0, 0, W, bar_h, fill=DARK)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # left-pad via paragraph indent
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(0.3))))
    return r


def textbox(slide, l, t, w, h, text="", size=Pt(18), bold=False,
            color=MID, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Simple single-run textbox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bullet_box(slide, l, t, w, h, items, size=Pt(17), title=None,
               title_size=Pt(19), color=MID, title_color=DARK,
               indent=Inches(0.25), spacing=Pt(6)):
    """Textbox with optional bold title then bullet items."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    para_idx = 0
    if title:
        p = tf.paragraphs[para_idx] if para_idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = title_size
        run.font.bold = True
        run.font.color.rgb = title_color
        _set_space_after(p, Pt(4))
        para_idx += 1

    for item in items:
        p = tf.add_paragraph() if (para_idx > 0 or title) else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_para_indent(p, indent)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = size
        run.font.color.rgb = color
        _set_space_after(p, spacing)
        para_idx += 1
    return txb


def _set_space_after(para, pt):
    pPr = para._p.get_or_add_pPr()
    spcAft = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft')
    spcPts = etree.SubElement(spcAft, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
    spcPts.set('val', str(int(pt.pt * 100)))


def _set_para_indent(para, indent):
    pPr = para._p.get_or_add_pPr()
    pPr.set('marL', str(int(indent)))
    pPr.set('indent', str(int(-Inches(0.15))))


def label(slide, l, t, w, h, text, size=Pt(14), fill=C1,
          text_color=WHITE, bold=True, radius=True):
    """Coloured label chip."""
    r = rect(slide, l, t, w, h, fill=fill)
    if radius:
        sp_pr = r._element.spPr
        prstGeom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is None:
            prstGeom = etree.SubElement(sp_pr,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            prstGeom.set('prst', 'roundRect')
            etree.SubElement(prstGeom,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = text_color
    return r


def divider(slide, y, color=GRAY, width=Pt(1)):
    """Horizontal rule."""
    line = slide.shapes.add_connector(1, Inches(0.4), y, W - Inches(0.4), y)
    line.line.color.rgb = color
    line.line.width = width


def callout_box(slide, l, t, w, h, title, body, t_color=DARK,
                fill=RGBColor(0xE8, 0xF4, 0xE8), border=GREEN,
                t_size=Pt(16), b_size=Pt(15)):
    """Coloured callout with bold title + body text."""
    r = rect(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(1.5))
    tf = r.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = t_size
    run.font.bold = True
    run.font.color.rgb = t_color
    _set_space_after(p, Pt(4))
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = body
    run2.font.size = b_size
    run2.font.color.rgb = MID
    return r


def alert_box(slide, l, t, w, h, title, body,
              fill=RGBColor(0xFF, 0xED, 0xED), border=ACCENT,
              t_size=Pt(16), b_size=Pt(15)):
    return callout_box(slide, l, t, w, h, title, body,
                       t_color=ACCENT, fill=fill, border=border,
                       t_size=t_size, b_size=b_size)


def note_box(slide, l, t, w, h, title, body,
             fill=RGBColor(0xE8, 0xF0, 0xFF), border=C1,
             t_size=Pt(16), b_size=Pt(15)):
    return callout_box(slide, l, t, w, h, title, body,
                       t_color=C1, fill=fill, border=border,
                       t_size=t_size, b_size=b_size)


def add_table(slide, l, t, w, rows_data, col_widths,
              header_fill=DARK, header_text=WHITE,
              header_size=Pt(13), body_size=Pt(13),
              row_height=Inches(0.38), alt_fill=RGBColor(0xF5, 0xF7, 0xFF),
              highlight_rows=None, highlight_fill=YELLOW):
    """Draw a table as individual rect+text shapes (avoids pptx table quirks)."""
    n_cols = len(col_widths)
    n_rows = len(rows_data)
    y = t
    for r_idx, row in enumerate(rows_data):
        x = l
        is_header = (r_idx == 0)
        is_divider = isinstance(row, str) and row == "---"
        if is_divider:
            divider(slide, y + row_height // 2)
            y += Inches(0.08)
            x = l
            continue
        for c_idx, cell in enumerate(row):
            cw = col_widths[c_idx]
            bold_cell = is_header or (c_idx == 0 and r_idx > 0)
            if is_header:
                fill = header_fill
                tc = header_text
            elif highlight_rows and r_idx in highlight_rows:
                fill = highlight_fill
                tc = DARK
            elif r_idx % 2 == 0:
                fill = alt_fill
                tc = MID
            else:
                fill = WHITE
                tc = MID
            r_shape = rect(slide, x, y, cw, row_height, fill=fill,
                           line=RGBColor(0xCC, 0xCC, 0xDD), line_w=Pt(0.5))
            tf = r_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell)
            run.font.size = body_size if not is_header else header_size
            run.font.bold = bold_cell
            run.font.color.rgb = tc
            x += cw
        y += row_height
    return y  # bottom y


def arrow(slide, x1, y1, x2, y2, color=MID, width=Pt(1.5)):
    """Add a line connector (arrow head at end)."""
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)

# accent stripe
rect(s, 0, Inches(2.8), W, Inches(0.08), fill=C2)

textbox(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.5),
        "LectureBench", size=Pt(52), bold=True, color=WHITE, wrap=False)
textbox(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.55),
        "Long-Lecture Video RAG Benchmark", size=Pt(26), color=C2, wrap=False)
textbox(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.5),
        "Temporally Grounded QA over Academic Lectures",
        size=Pt(20), color=RGBColor(0xAA, 0xBB, 0xDD), italic=True)

rect(s, Inches(0.8), Inches(4.1), Inches(6), Inches(1.8),
     fill=RGBColor(0x28, 0x28, 0x48), line=RGBColor(0x55, 0x66, 0x88))
textbox(s, Inches(1.0), Inches(4.25), Inches(5.5), Inches(0.4),
        "Nasim Faridnia", size=Pt(18), bold=True, color=WHITE)
textbox(s, Inches(1.0), Inches(4.65), Inches(5.5), Inches(0.55),
        "Department of Computer Science\nUniversity of Texas at San Antonio",
        size=Pt(14), color=RGBColor(0xAA, 0xBB, 0xDD))
textbox(s, Inches(1.0), Inches(5.35), Inches(5.5), Inches(0.35),
        "BMVC 2026", size=Pt(14), color=C2)

textbox(s, Inches(8.5), Inches(4.2), Inches(4.3), Inches(1.6),
        "github.com/NassimF/VQA-Benchmark\n\n60 lectures  |  810 QA pairs\n62.6 hours  |  CC BY-NC-SA 4.0",
        size=Pt(14), color=RGBColor(0x99, 0xAA, 0xCC))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "The Problem: Long Lectures Are Underserved by Existing Benchmarks")

# left column
bullet_box(s, Inches(0.4), Inches(1.15), Inches(5.9), Inches(2.3),
    title="Short-clip benchmarks dominate VideoQA",
    items=[
        "ActivityNet-QA, NExT-QA: avg < 3 min per clip",
        "EgoSchema: 3-minute egocentric clips",
        "Optimized for moment classification, not evidence retrieval",
    ], size=Pt(16), title_size=Pt(18))

bullet_box(s, Inches(0.4), Inches(3.55), Inches(5.9), Inches(2.3),
    title="Lecture videos are fundamentally different",
    items=[
        "47–89 min per lecture",
        "Answers may span segments 60+ minutes apart",
        "Critical evidence lives on slides/whiteboards — absent from spoken transcript",
    ], size=Pt(16), title_size=Pt(18))

# right column
bullet_box(s, Inches(6.7), Inches(1.15), Inches(6.2), Inches(2.3),
    title="Current lecture benchmarks fall short",
    items=[
        "EduVidQA: fixed 4-min context window — retrieval problem sidestepped entirely",
        "No benchmark evaluates temporal grounding in lecture RAG",
        "No multi-hop evidence across non-adjacent segments",
    ], size=Pt(16), title_size=Pt(18))

alert_box(s, Inches(6.7), Inches(3.7), Inches(6.2), Inches(1.9),
    "The Gap",
    "No benchmark evaluates whether a RAG system retrieves the right moment "
    "from a long lecture — not merely a semantically similar one.",
    t_size=Pt(17), b_size=Pt(16))

divider(s, Inches(1.1))
divider(s, Inches(3.45))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Prior Work Comparison Table
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Comparison with Related Benchmarks")

headers = ["Dataset", "Domain", "#Videos", "#QA Pairs", "Avg Duration", "Multi-hop", "Temporal+Visual"]
rows = [
    headers,
    ["NExT-QA",       "General",    "5,440", "52,044", "44 sec",     "✗", "✗"],
    ["EgoSchema",     "Egocentric", "—",     "5,031",  "3 min",      "✗", "✗"],
    ["NA-VQA",        "Movies",     "88",    "4,442",  "80–200 min", "implicit", "✗"],
    ["LongVidSearch", "Mixed",      "447",   "3,000",  "26 min",     "✓", "✗"],
    ["EduVidQA",      "Lectures",   "296",   "5,252",  "68 min",     "✗", "✓"],
    ["LectureBench (ours)", "Lectures", "60", "810",   "63 min",     "✓", "✓"],
]
cw = [Inches(2.3), Inches(1.4), Inches(1.1), Inches(1.3),
      Inches(1.5), Inches(1.4), Inches(1.9)]
add_table(s, Inches(0.35), Inches(1.1), W - Inches(0.7), rows, cw,
          highlight_rows={6}, row_height=Inches(0.44), body_size=Pt(14))

# gap callouts at bottom
gap_y = Inches(5.65)
for i, (txt, c) in enumerate([
    ("Gap 1: No benchmark is\nmulti-hop AND lecture-domain", C1),
    ("Gap 2: No benchmark has\ntimestamp spans AND visual split", C2),
    ("Gap 3: EduVidQA sidesteps\nretrieval with 4-min windows", ACCENT),
]):
    xl = Inches(0.4) + i * Inches(4.3)
    r = rect(s, xl, gap_y, Inches(4.0), Inches(0.95),
             fill=RGBColor(0xEE, 0xF2, 0xFF) if c == C1
                  else (RGBColor(0xFF, 0xF0, 0xE8) if c == C2
                        else RGBColor(0xFF, 0xED, 0xED)),
             line=c, line_w=Pt(1.5))
    tf = r.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(14)
    run.font.color.rgb = c
    run.font.bold = True


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Contributions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Our Contributions")

contribs = [
    (C1,    "1",
     "Corpus",
     "60 CC-licensed academic lecture videos (62.6 hours) across 15 courses "
     "from MIT OCW, NYU Deep Learning, and Stanford SEE."),
    (C2,    "2",
     "Benchmark",
     "810 LLM-reviewed QA pairs with precise temporal ground-truth spans, "
     "typed by multi-hop-visual, visual, multi-hop, text, and unanswerable categories."),
    (GREEN, "3",
     "Open-source RAG Pipeline",
     "Two retrieval configurations (transcript-only vs transcript+frames) with a "
     "full evaluation suite: temporal IoU, Hit Rate@k, LLM-judge, citation accuracy."),
    (ACCENT,"4",
     "Empirical Finding",
     "Visual frame captions yield +65% tIoU on visual-dependent questions "
     "and are essentially unchanged (–3%) on text-only questions — "
     "precisely isolating the multimodal contribution."),
]

for i, (color, num, ctitle, body) in enumerate(contribs):
    y = Inches(1.2) + i * Inches(1.42)
    # number circle
    label(s, Inches(0.35), y + Inches(0.08), Inches(0.55), Inches(0.55),
          num, size=Pt(20), fill=color)
    # title
    textbox(s, Inches(1.1), y + Inches(0.03), Inches(11.6), Inches(0.42),
            ctitle, size=Pt(19), bold=True, color=color)
    # body
    textbox(s, Inches(1.1), y + Inches(0.48), Inches(11.6), Inches(0.72),
            body, size=Pt(16), color=MID)
    divider(s, y + Inches(1.32), color=RGBColor(0xDD, 0xDD, 0xEE))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LectureBench at a Glance
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "LectureBench at a Glance")

# left: corpus stats table
stats = [
    ["Statistic", "Value"],
    ["Lectures", "60"],
    ["Sources", "MIT OCW, NYU DL, Stanford SEE"],
    ["Total duration", "62.6 hours"],
    ["Duration range", "47–89 min (avg 63 min)"],
    ["Transcript chunks (45s/35s)", "7,382"],
    ["Frames captioned (every 15s)", "17,180"],
    ["QA pairs (accepted)", "810 (avg 13.5/lecture)"],
    ["Visual-dependent pairs", "56.2%  (455/810)"],
]
add_table(s, Inches(0.35), Inches(1.1), Inches(6.4), stats,
          [Inches(3.4), Inches(3.0)], row_height=Inches(0.4), body_size=Pt(14))

# right: QA type breakdown
qt_rows = [
    ["Question Type", "Count", "Role"],
    ["Multi-hop visual", "320", "Primary"],
    ["Single-hop visual", "135", "Primary"],
    ["Multi-hop textual", "140", "Control"],
    ["Single-hop textual", "103", "Control"],
    ["Unanswerable", "112", "Required"],
    ["Total", "810", ""],
]
add_table(s, Inches(7.0), Inches(1.1), Inches(5.9), qt_rows,
          [Inches(2.9), Inches(1.2), Inches(1.8)],
          row_height=Inches(0.44), body_size=Pt(15), highlight_rows={6})

# figure placeholder
rect(s, Inches(7.0), Inches(4.3), Inches(5.9), Inches(2.6),
     fill=RGBColor(0xE8, 0xEC, 0xF8), line=C1, line_w=Pt(1))
textbox(s, Inches(7.0), Inches(5.35), Inches(5.9), Inches(0.5),
        "[ Insert Figure 6 — QA type distribution bar chart ]",
        size=Pt(13), color=C1, align=PP_ALIGN.CENTER, italic=True)

textbox(s, Inches(0.35), Inches(6.95), Inches(7.0), Inches(0.35),
        "Dataset released under CC BY-NC-SA 4.0  |  github.com/NassimF/VQA-Benchmark",
        size=Pt(12), color=RGBColor(0x77, 0x77, 0x99), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline Overview (Figure 2 recreation)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "End-to-End Pipeline  (Figure 2)")

# ── phase backgrounds ────────────────────────────────────────────────────────
rect(s, Inches(0.25), Inches(1.05), Inches(7.0), Inches(5.65),
     fill=RGBColor(0xEE, 0xF2, 0xF8), line=RGBColor(0xBB, 0xCC, 0xDD))
rect(s, Inches(7.5), Inches(1.05), Inches(5.5), Inches(5.65),
     fill=RGBColor(0xFE, 0xF6, 0xEF), line=RGBColor(0xDD, 0xCC, 0xBB))

textbox(s, Inches(0.35), Inches(1.08), Inches(3.5), Inches(0.38),
        "Offline Indexing", size=Pt(15), bold=True, color=RGBColor(0x44, 0x44, 0x66))
textbox(s, Inches(7.6), Inches(1.08), Inches(4.0), Inches(0.38),
        "Online Query & Evaluation", size=Pt(15), bold=True, color=RGBColor(0x44, 0x44, 0x66))

# ── shared input row (top) ───────────────────────────────────────────────────
bh, bw = Inches(0.55), Inches(1.5)   # box height / width
row1_y = Inches(1.65)

def pipe_box(slide, x, y, w, h, text, fill, tcolor=WHITE):
    r = rect(slide, x, y, w, h, fill=fill, line=RGBColor(0x88,0x88,0x88), line_w=Pt(0.75))
    tf = r.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = tcolor
    return r

def h_arrow(slide, x1, x2, y, color=MID):
    arrow(slide, x1, y, x2, y, color=color, width=Pt(1.5))

NGRAY = RGBColor(0x77, 0x77, 0x88)

# Row 1 boxes
pipe_box(s, Inches(0.4),  row1_y, bw, bh, "Lecture\nVideo",       NGRAY)
pipe_box(s, Inches(2.1),  row1_y, bw, bh, "VTT\nTranscript",      NGRAY)
pipe_box(s, Inches(3.8),  row1_y, bw, bh, "Chunker\n45s/35s",     NGRAY)
pipe_box(s, Inches(5.5),  row1_y, Inches(1.7), bh,
         "Config 1\nTranscript Only\n(ChromaDB)", C1)

# Row 1 arrows
h_arrow(s, Inches(1.9),  Inches(2.1),  row1_y + bh/2, NGRAY)
h_arrow(s, Inches(3.6),  Inches(3.8),  row1_y + bh/2, NGRAY)
h_arrow(s, Inches(5.3),  Inches(5.5),  row1_y + bh/2, C1)

# Row 2 boxes (Config 2 branch)
row2_y = Inches(3.1)
pipe_box(s, Inches(3.8),  row2_y, bw, bh,
         "Frame Caps\nQwen2-VL-7B", C2)
pipe_box(s, Inches(5.5),  row2_y, Inches(1.7), bh,
         "Config 2\nText + Captions\n(ChromaDB)", C2)

# Video → Frame Caps arrow (down then right)
arrow(s, Inches(0.4)+bw/2, row1_y+bh, Inches(0.4)+bw/2, row2_y+bh/2, C2, Pt(1.5))
h_arrow(s, Inches(0.4)+bw/2, Inches(3.8), row2_y+bh/2, C2)
h_arrow(s, Inches(5.3),  Inches(5.5),  row2_y+bh/2, C2)

# Shared online boxes
on_x1, on_x2, on_x3 = Inches(7.6), Inches(9.3), Inches(11.0)
pipe_box(s, on_x1, row1_y, bw, bh, "Retriever\ntop-k cosine", NGRAY)
pipe_box(s, on_x2, row1_y, bw, bh, "Generator\nGPT-4o-mini",  NGRAY)
pipe_box(s, on_x3, row1_y, bw, bh, "Answer +\nCitations",     NGRAY)

# C1 → retriever
arrow(s, Inches(5.5)+Inches(1.7), row1_y+bh/2, on_x1, row1_y+bh/2, C1, Pt(1.5))
# C2 → retriever (from below)
arrow(s, Inches(5.5)+Inches(1.7), row2_y+bh/2,
         on_x1 + bw/2,             row2_y+bh/2, C2, Pt(1.5))
arrow(s, on_x1 + bw/2,            row2_y+bh/2,
         on_x1 + bw/2,             row1_y+bh,   C2, Pt(1.5))

h_arrow(s, on_x1+bw, on_x2,     row1_y+bh/2, NGRAY)
h_arrow(s, on_x2+bw, on_x3,     row1_y+bh/2, NGRAY)

# Evaluation row
ev_y = Inches(3.1)
pipe_box(s, on_x1, ev_y, bw, bh, "GT Spans\n(Benchmark)", GREEN)
pipe_box(s, on_x2, ev_y, bw, bh, "Evaluator", GREEN)
pipe_box(s, on_x3, ev_y, bw, bh,
         "tIoU · HR@k\nJudge · Cit.Acc", GREEN)

# Answer → Evaluator
arrow(s, on_x3+bw/2, row1_y+bh, on_x2+bw*0.75, ev_y, GREEN, Pt(1.5))
h_arrow(s, on_x1+bw, on_x2, ev_y+bh/2, GREEN)
h_arrow(s, on_x2+bw, on_x3, ev_y+bh/2, GREEN)

# Caption below pipeline
textbox(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(0.8),
        "Both configs share the same retriever, generator, and evaluator — "
        "the only variable is chunk content (transcript text vs. transcript + frame captions).",
        size=Pt(15), color=MID, italic=True)

# legend
for lx, lc, lt in [(Inches(0.4), C1, "Config 1 (transcript only)"),
                    (Inches(3.5), C2, "Config 2 (transcript + frame captions)"),
                    (Inches(7.5), GREEN, "Evaluation")]:
    rect(s, lx, Inches(5.5), Inches(0.28), Inches(0.22), fill=lc)
    textbox(s, lx + Inches(0.35), Inches(5.47), Inches(2.8), Inches(0.3),
            lt, size=Pt(13), color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Chunking & Frame Captioning
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Pipeline: Temporal Chunking & Frame Captioning")

# left
bullet_box(s, Inches(0.4), Inches(1.2), Inches(6.0), Inches(2.7),
    title="Temporal Chunking",
    items=[
        "Window: 45 seconds  |  Stride: 35 seconds (10s overlap)",
        "45s ≈ one complete concept explanation per lecture",
        "Overlap ensures evidence near chunk boundaries appears in both neighbours",
        "7,382 chunks across 60 lectures",
        "No shot boundaries in lectures — fixed windows are the only practical choice",
    ], size=Pt(16), title_size=Pt(19))

note_box(s, Inches(0.4), Inches(4.0), Inches(6.0), Inches(1.2),
    "Why not shot-boundary chunking?",
    "Lecture recordings have no camera cuts. Fixed-width sliding windows keep the "
    "pipeline simple and the comparison interpretable.",
    t_size=Pt(16), b_size=Pt(15))

# right
bullet_box(s, Inches(6.8), Inches(1.2), Inches(6.1), Inches(3.0),
    title="Frame Captioning  (Config 2 only)",
    items=[
        "1 frame extracted every 15s via ffmpeg",
        "Captioned by Qwen2-VL-7B-Instruct (bfloat16, single A100)",
        "Prompt: describe all visible text, equations, diagrams in detail",
        "Caption appended as [frame caption: …] to chunk text",
        "17,180 frames total  |  ~3 frames per chunk",
        "Cost: ~2–8 GPU-hours for all 60 lectures",
    ], size=Pt(16), title_size=Pt(19))

callout_box(s, Inches(6.8), Inches(4.25), Inches(6.1), Inches(1.1),
    "Key Design Choice",
    "Appending captions as plain text keeps the embedding model fixed across both "
    "configs — performance differences are due to chunk content alone.",
    t_size=Pt(16), b_size=Pt(15))

divider(s, Inches(1.15))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Retrieval & Generation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Pipeline: Retrieval & Answer Generation")

bullet_box(s, Inches(0.4), Inches(1.2), Inches(6.1), Inches(3.2),
    title="Retrieval",
    items=[
        "Embedding: all-MiniLM-L6-v2  (384-dim, Sentence-BERT)",
        "Vector store: ChromaDB — two persistent collections",
        "   • transcript_only",
        "   • transcript_plus_frames",
        "Query-time: top-k cosine similarity,  k ∈ {1, 3, 5}",
        "Same model + store for both configs",
    ], size=Pt(16), title_size=Pt(19))

note_box(s, Inches(0.4), Inches(4.55), Inches(6.1), Inches(1.1),
    "Controlled comparison",
    "Fixing the embedding model means any difference in results is solely "
    "attributable to what text was embedded — not to retrieval capacity.",
    t_size=Pt(16), b_size=Pt(15))

bullet_box(s, Inches(6.8), Inches(1.2), Inches(6.1), Inches(3.2),
    title="Answer Generation",
    items=[
        "Generator: GPT-4o-mini at temperature 0.1",
        "Input: retrieved chunks as numbered excerpts",
        "Output: structured JSON  →  answer  +  citations list",
        "Predicted span = union of all cited chunk time windows",
        "Citations rendered as [video_id @ mm:ss to mm:ss]",
    ], size=Pt(16), title_size=Pt(19))

alert_box(s, Inches(6.8), Inches(4.55), Inches(6.1), Inches(1.1),
    "Why a weaker generator?",
    "Using GPT-4o-mini isolates retrieval failures that a stronger model might mask "
    "by drawing on world knowledge outside the retrieved chunks.",
    t_size=Pt(16), b_size=Pt(15))

divider(s, Inches(1.15))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Benchmark Construction
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Benchmark Construction: QA Generation & LLM Review")

bullet_box(s, Inches(0.4), Inches(1.2), Inches(6.1), Inches(3.0),
    title="QA Generation",
    items=[
        "Claude Sonnet 4.6 generates 15 draft QA pairs per lecture",
        "Input: augmented chunks (transcript + frame captions)",
        "5 types targeted: multi-hop-visual (7), visual (2), multi-hop (3), text (2), unanswerable (1)",
        "Multi-hop: ≥2 non-adjacent spans, gap ≥ 70s required",
        "Total API cost: ~$8 for all 60 lectures",
    ], size=Pt(16), title_size=Pt(19))

note_box(s, Inches(0.4), Inches(4.35), Inches(6.1), Inches(1.2),
    "Generator ≠ Evaluator (intentional)",
    "QA generation uses Claude Sonnet 4.6; answer generation at eval time uses "
    "GPT-4o-mini. Using the same model for both would be circular — the generator "
    "would reproduce its own phrasing regardless of retrieval quality.",
    t_size=Pt(15), b_size=Pt(14))

# right: review pipeline as steps
bullet_box(s, Inches(6.8), Inches(1.2), Inches(6.1), Inches(3.6),
    title="LLM Review (two-pass filter)",
    items=[
        "Pass 1 — Claude Sonnet 4.6 — three-criterion binary rubric:",
        "   C1: answer correctness vs. chunk text",
        "   C2: span plausibility (≥70s gap between hops)",
        "   C3: question type accuracy",
        "Pass 2 — Cross-family check: pairs passing C1 re-evaluated by GPT-4o",
        "   Disagreement between models triggers rejection",
        "Atomic fact decomposition (FActScore) before C1 check",
        "Acceptance rate: ~90%  (810 accepted / 900 target)",
    ], size=Pt(15), title_size=Pt(19))

divider(s, Inches(1.15))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Example QA Pair
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Example QA Pair: Multi-hop Visual  (mit_6046_lec01_q017)")

# fbox background
rect(s, Inches(0.35), Inches(1.15), Inches(12.6), Inches(4.5),
     fill=WHITE, line=RGBColor(0x88, 0x88, 0xAA), line_w=Pt(1))

textbox(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.45),
        "Question (multi-hop-visual):", size=Pt(16), bold=True, color=DARK)
textbox(s, Inches(0.55), Inches(1.68), Inches(12.2), Inches(0.55),
        "What does the blackboard show as the theme of today's lecture, and how does the "
        "instructor connect this theme to the interval scheduling problems he discusses?",
        size=Pt(15), color=MID)

textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(0.38),
        "Ground-truth spans:", size=Pt(16), bold=True, color=DARK)

# Hop 1
rect(s, Inches(0.55), Inches(2.72), Inches(12.1), Inches(0.72),
     fill=RGBColor(0xEE, 0xF4, 0xFF), line=C1, line_w=Pt(1))
textbox(s, Inches(0.7), Inches(2.76), Inches(11.8), Inches(0.62),
        "Hop 1  [10:20–11:15]  (visual) — frame caption: Blackboard shows "
        "\"Theme of today's lecture — Very similar problems can have very different complexity.\"",
        size=Pt(14), color=C1)

# Hop 2
rect(s, Inches(0.55), Inches(3.5), Inches(12.1), Inches(0.72),
     fill=RGBColor(0xFF, 0xF4, 0xEE), line=C2, line_w=Pt(1))
textbox(s, Inches(0.7), Inches(3.54), Inches(11.8), Inches(0.62),
        "Hop 2  [79:35–80:40]  (transcript) — \"…interval scheduling is greedy O(n); "
        "weighted interval scheduling requires O(n²) DP; non-identical machines is NP-complete…\"",
        size=Pt(14), color=C2)

textbox(s, Inches(0.55), Inches(4.28), Inches(12.2), Inches(0.35),
        "Span gap:  4775 − 675 = 4100 s  (68 min)  >>  70 s required",
        size=Pt(15), bold=True, color=RGBColor(0x33, 0x88, 0x33))

# figure placeholder
rect(s, Inches(0.35), Inches(5.75), Inches(12.6), Inches(1.3),
     fill=RGBColor(0xE8, 0xEC, 0xF8), line=C1, line_w=Pt(1))
textbox(s, Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.45),
        "[ Insert Figure 3 — timeline diagram showing Hop 1 at 10:20 and Hop 2 at 79:35 with 68-min gap annotation ]",
        size=Pt(13), color=C1, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Evaluation Metrics
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Evaluation Metrics")

metrics = [
    (C1,    "Temporal Grounding",
     "Mean Temporal IoU — overlap between predicted span (union of cited chunks) and ground-truth span.\n"
     "IoU@0.3 / IoU@0.5 — fraction of pairs exceeding threshold."),
    (C2,    "Retrieval Coverage",
     "Hit Rate@k  (k ∈ {1, 3, 5}) — fraction of questions where the ground-truth span is covered by "
     "at least one of the top-k retrieved chunks at IoU ≥ 0.3."),
    (GREEN, "Answer Quality",
     "LLM-Judge Score (1–5) — GPT-4o scores the generated answer against the reference; "
     "85% agreement with human experts (Zheng et al. 2023)."),
    (ACCENT,"Citation Accuracy",
     "Fraction of cited time spans that overlap the correct ground-truth evidence window. "
     "Measures whether the system knows what evidence it used."),
]
for i, (color, mname, mbody) in enumerate(metrics):
    y = Inches(1.25) + i * Inches(1.42)
    label(s, Inches(0.35), y, Inches(2.9), Inches(0.5),
          mname, size=Pt(15), fill=color)
    textbox(s, Inches(3.5), y, Inches(9.5), Inches(0.95),
            mbody, size=Pt(16), color=MID)
    divider(s, y + Inches(1.25), color=RGBColor(0xDD,0xDD,0xEE))

callout_box(s, Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.42),
    "Primary metric: mean tIoU — ",
    "At the ±25.8s span annotation error, tIoU@0.3 remains valid (IoU ≈ 0.4–0.6 at that error level).",
    t_size=Pt(14), b_size=Pt(14))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Single-hop results
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Results: Single-Hop Questions  (n = 238)")

rows = [
    ["Metric", "Config 1\n(Transcript-Only)", "Config 2\n(+Frames)", "Change"],
    ["Mean Temporal IoU",  "0.175", "0.281", "+60%"],
    ["IoU@0.3",           "0.248", "0.387", "+56%"],
    ["IoU@0.5",           "0.164", "0.290", "+77%"],
    ["Hit Rate@1",        "0.218", "0.244", "+12%"],
    ["Hit Rate@3",        "0.382", "0.420", "+10%"],
    ["Hit Rate@5",        "0.445", "0.529", "+19%"],
    ["LLM-Judge (1–5)",   "2.68",  "3.53",  "+32%"],
    ["Citation Accuracy", "0.269", "0.433", "+61%"],
]
add_table(s, Inches(0.35), Inches(1.15), Inches(7.2), rows,
          [Inches(2.5), Inches(1.9), Inches(1.9), Inches(0.9)],
          row_height=Inches(0.45), body_size=Pt(15), highlight_rows={1,8})

bullet_box(s, Inches(7.85), Inches(1.25), Inches(5.1), Inches(4.5),
    title="Key Observations",
    items=[
        "Config 2 wins on all 8 metrics",
        "Largest gains: tIoU +60%, citation accuracy +61%",
        "LLM-judge +32% — better retrieval → more accurate answers",
        "Hit Rate gains are modest (+12% at k=1): frame captions improve the quality of retrieved chunks, not their presence in top-k",
        "135/238 pairs (57%) are visual-dependent — these drive the large tIoU and citation gains",
    ], size=Pt(16), title_size=Pt(19))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Multi-hop results
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Results: Multi-Hop Questions  (n = 460)")

rows = [
    ["Metric", "Config 1\n(Transcript-Only)", "Config 2\n(+Frames)", "Change"],
    ["Mean Temporal IoU",  "0.180", "0.204", "+13%"],
    ["IoU@0.3",           "0.274", "0.291", "+6%"],
    ["IoU@0.5",           "0.076", "0.072", "–5%*"],
    ["Hit Rate@1",        "0.326", "0.341", "+5%"],
    ["Hit Rate@3",        "0.563", "0.620", "+10%"],
    ["Hit Rate@5",        "0.676", "0.722", "+7%"],
    ["LLM-Judge (1–5)",   "2.95",  "3.45",  "+17%"],
    ["Citation Accuracy", "0.273", "0.333", "+22%"],
]
add_table(s, Inches(0.35), Inches(1.15), Inches(7.2), rows,
          [Inches(2.5), Inches(1.9), Inches(1.9), Inches(0.9)],
          row_height=Inches(0.45), body_size=Pt(15))

bullet_box(s, Inches(7.85), Inches(1.25), Inches(5.1), Inches(3.5),
    title="Key Observations",
    items=[
        "Config 2 wins 7/8 metrics",
        "Smaller gains than single-hop (+13% tIoU vs +60%)",
        "Mixed composition: 320/460 visual pairs benefit; 140/460 text-only show slight regression",
        "Hit Rate@5 is 107% higher than Hit Rate@1 — multi-hop retrieval is k-sensitive",
    ], size=Pt(16), title_size=Pt(19))

textbox(s, Inches(7.85), Inches(4.85), Inches(5.1), Inches(0.9),
        "* IoU@0.5 footnote: frame-augmented chunks are longer, widening the union span "
        "and reducing tIoU even when the correct chunk is retrieved.",
        size=Pt(13), color=RGBColor(0x99, 0x55, 0x55), italic=True)

# figure placeholder
rect(s, Inches(7.85), Inches(5.85), Inches(5.1), Inches(1.35),
     fill=RGBColor(0xE8, 0xEC, 0xF8), line=C1, line_w=Pt(1))
textbox(s, Inches(7.85), Inches(6.3), Inches(5.1), Inches(0.4),
        "[ Insert Figure 4 — Hit Rate@k line plot ]",
        size=Pt(13), color=C1, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Key Finding: Visual vs Text split
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Key Finding: Visual vs. Text-Only Question Split")

rows = [
    ["Question Type", "Config 1\ntIoU", "Config 1\nJudge", "Config 1\nCit.Acc",
                      "Config 2\ntIoU", "Config 2\nJudge", "Config 2\nCit.Acc"],
    ["Visual-dep.  (n=455)", "0.129", "2.57", "0.198", "0.213", "3.52", "0.340"],
    ["Text-only    (n=243)", "0.271", "3.38", "0.409", "0.263", "3.41", "0.418"],
]
add_table(s, Inches(0.35), Inches(1.15), Inches(12.6), rows,
          [Inches(2.5), Inches(1.5), Inches(1.4), Inches(1.5),
           Inches(1.5), Inches(1.4), Inches(1.5)],
          row_height=Inches(0.55), body_size=Pt(15), highlight_rows={1})

# two callouts
callout_box(s, Inches(0.35), Inches(2.75), Inches(6.1), Inches(2.1),
    "Visual-dependent questions",
    "Frame captions yield:\n"
    "  +65% tIoU   (0.129 → 0.213)\n"
    "  +37% LLM-judge   (2.57 → 3.52)\n"
    "  +72% citation accuracy   (0.198 → 0.340)",
    fill=RGBColor(0xE0, 0xF4, 0xE0), border=GREEN,
    t_color=GREEN, t_size=Pt(17), b_size=Pt(16))

alert_box(s, Inches(6.85), Inches(2.75), Inches(6.1), Inches(2.1),
    "Text-only questions",
    "Essentially unchanged:\n"
    "  –3% tIoU   (0.271 → 0.263)\n"
    "  +1% LLM-judge   (3.38 → 3.41)\n"
    "  +2% citation accuracy   (0.409 → 0.418)",
    fill=RGBColor(0xFF, 0xF8, 0xF0), border=RGBColor(0xCC, 0x88, 0x44),
    t_size=Pt(17), b_size=Pt(16))

rect(s, Inches(0.35), Inches(5.1), Inches(12.6), Inches(0.7),
     fill=DARK, line=None)
textbox(s, Inches(0.35), Inches(5.18), Inches(12.6), Inches(0.55),
        "Frame captions add signal only where visual evidence is required — and introduce no noise elsewhere.",
        size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Failure Mode Analysis
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Failure Mode Analysis")

# 3 failure mode cards
for i, (color, title_, body_) in enumerate([
    (C1,    "1. Span Inflation",
     "Entire 45s chunk cited when relevant evidence occupies only a fraction → "
     "tIoU penalty even when correct chunk is retrieved."),
    (ACCENT,"2. Visual Blindness  (Config 1)",
     "Answer appears only on a lecture slide — absent from spoken transcript. "
     "Config 2 recovers by reading the frame caption.\n"
     "Mean tIoU:  0.129 → 0.213  (+65%)"),
    (C2,    "3. Partial Multi-hop",
     "Retriever surfaces one hop at rank 1 but misses the second. "
     "Hit Rate@5 is 107% higher than Hit Rate@1 on multi-hop questions."),
]):
    xl = Inches(0.35) + i * Inches(4.33)
    rect(s, xl, Inches(1.15), Inches(4.1), Inches(2.2),
         fill=RGBColor(0xF5,0xF7,0xFF) if color==C1
              else (RGBColor(0xFF,0xED,0xED) if color==ACCENT
                    else RGBColor(0xFF,0xF5,0xEE)),
         line=color, line_w=Pt(2))
    textbox(s, xl+Inches(0.12), Inches(1.22), Inches(3.86), Inches(0.45),
            title_, size=Pt(17), bold=True, color=color)
    textbox(s, xl+Inches(0.12), Inches(1.7), Inches(3.86), Inches(1.5),
            body_, size=Pt(14), color=MID)

# visual-blindness example box
rect(s, Inches(0.35), Inches(3.55), Inches(12.6), Inches(3.4),
     fill=WHITE, line=RGBColor(0x88,0x88,0xAA), line_w=Pt(1))

textbox(s, Inches(0.55), Inches(3.65), Inches(12.2), Inches(0.38),
        "Visual-blindness case  (nyu_dl_week10_q008)", size=Pt(15), bold=True, color=DARK)
textbox(s, Inches(0.55), Inches(4.05), Inches(12.2), Inches(0.38),
        "Q: What are the three bullet points on the 'Why self-supervision?' slide?",
        size=Pt(14), color=MID)

rect(s, Inches(0.55), Inches(4.48), Inches(5.8), Inches(1.55),
     fill=RGBColor(0xFF,0xED,0xED), line=ACCENT, line_w=Pt(1))
textbox(s, Inches(0.68), Inches(4.55), Inches(5.5), Inches(0.35),
        "Config 1 — transcript only", size=Pt(13), bold=True, color=ACCENT)
textbox(s, Inches(0.68), Inches(4.93), Inches(5.5), Inches(0.72),
        "Answer: \"The provided excerpts do not contain enough information.\"\n"
        "LLM-judge: 1.5 / 5", size=Pt(13), color=MID)

rect(s, Inches(6.7), Inches(4.48), Inches(6.15), Inches(1.55),
     fill=RGBColor(0xE0,0xF4,0xE0), line=GREEN, line_w=Pt(1))
textbox(s, Inches(6.83), Inches(4.55), Inches(5.85), Inches(0.35),
        "Config 2 — same chunk + frame caption", size=Pt(13), bold=True, color=GREEN)
textbox(s, Inches(6.83), Inches(4.93), Inches(5.85), Inches(0.72),
        "Frame caption: Slide 'Why self-supervision?' — (1) Helps us learn using observations; "
        "(2) Does not require exhaustive annotation; (3) Leverage multiple modalities.\n"
        "LLM-judge: 5.0 / 5", size=Pt(13), color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Hit Rate @ k
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Hit Rate Analysis: k-Sensitivity of Multi-hop Retrieval")

# figure placeholder left
rect(s, Inches(0.35), Inches(1.15), Inches(6.5), Inches(5.55),
     fill=RGBColor(0xE8,0xEC,0xF8), line=C1, line_w=Pt(1))
textbox(s, Inches(0.35), Inches(3.7), Inches(6.5), Inches(0.45),
        "[ Insert Figure 4 — fig_hr_at_k.pdf ]",
        size=Pt(14), color=C1, align=PP_ALIGN.CENTER, italic=True)

# right observations
bullet_box(s, Inches(7.1), Inches(1.25), Inches(5.9), Inches(4.0),
    title="Reading the plot",
    items=[
        "Solid lines = single-hop;  dashed = multi-hop",
        "Config 2 (+Frames) is consistently above Config 1 at all k",
        "Single-hop curves plateau early — first hop dominates at low k",
    ], size=Pt(16), title_size=Pt(18))

bullet_box(s, Inches(7.1), Inches(3.5), Inches(5.9), Inches(2.7),
    title="Multi-hop curves: steep rise through k=5",
    items=[
        "Config 1:  HR@1 = 0.326  →  HR@5 = 0.676  (+107%)",
        "Config 2:  HR@1 = 0.341  →  HR@5 = 0.722  (+112%)",
        "The 2nd hop lands at rank 3–5, not rank 1",
        "Larger k budgets would directly improve multi-hop performance",
    ], size=Pt(16), title_size=Pt(18))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Summary of Results
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Summary of Results")

rows = [
    ["Metric", "Config 1", "Config 2", "Δ", "Where the gain comes from"],
    ["Mean tIoU (all)",       "0.154", "0.198", "+29%", "Visual +65%, text 0%"],
    ["IoU@0.3",               "0.228", "0.279", "+22%", "—"],
    ["Hit Rate@5",            "0.515", "0.565", "+10%", "—"],
    ["LLM-Judge (1–5)",       "3.03",  "3.56",  "+17%", "—"],
    ["Citation Accuracy",     "0.234", "0.316", "+35%", "—"],
    ["tIoU — visual (n=455)", "0.129", "0.213", "+65%", "Frame caps expose slide text"],
    ["tIoU — text   (n=243)", "0.271", "0.263", "–3%",  "No visual evidence needed"],
]
add_table(s, Inches(0.35), Inches(1.15), Inches(12.6), rows,
          [Inches(2.7), Inches(1.2), Inches(1.2), Inches(0.9), Inches(6.6)],
          row_height=Inches(0.48), body_size=Pt(14),
          highlight_rows={6}, header_size=Pt(14))

callout_box(s, Inches(0.35), Inches(5.85), Inches(6.1), Inches(1.0),
    "What works",
    "Frame captions close the visual-evidence gap precisely: +65% tIoU and +72% citation accuracy on visual-dependent questions.",
    t_size=Pt(16), b_size=Pt(15))

alert_box(s, Inches(6.85), Inches(5.85), Inches(6.1), Inches(1.0),
    "Precise multimodal augmentation",
    "No degradation on transcript-only questions (–3% tIoU, within noise). "
    "Frame captions add signal exactly where they should.",
    t_size=Pt(16), b_size=Pt(15))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Limitations
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Limitations")

limits = [
    "English-only corpus — limited to YouTube auto-captions; occasional transcription errors in technical terms.",
    "LLM-only QA review — ground-truth spans estimated from chunk boundaries (±15–30s), not human-tightened.",
    "Frame captioner not fine-tuned — Qwen2-VL-7B occasionally misreads small mathematical notation on slides.",
    "Scale — 810 pairs across 60 lectures is smaller than existing benchmarks (EduVidQA: 5,252; LongVidSearch: 3,000). Justified by higher per-pair curation cost (temporal grounding + multi-hop validation + cross-family review).",
    "Single embedding model — results specific to all-MiniLM-L6-v2; stronger embeddings may compress the gap.",
]
for i, lim in enumerate(limits):
    y = Inches(1.25) + i * Inches(1.1)
    label(s, Inches(0.35), y, Inches(0.55), Inches(0.45),
          str(i+1), size=Pt(16), fill=RGBColor(0x99,0x99,0xBB))
    textbox(s, Inches(1.1), y, Inches(11.8), Inches(0.95),
            lim, size=Pt(16), color=MID)
    if i < 4:
        divider(s, y + Inches(1.0), color=RGBColor(0xDD,0xDD,0xEE))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Future Work
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Future Work")

fw = [
    (C1,    "Non-English lectures",
     "Extend the corpus to multilingual content and evaluate ASR quality as an independent variable."),
    (C2,    "Semantic chunking",
     "Replace sliding-window with topic-boundary detection to produce semantically coherent chunks."),
    (GREEN, "Joint text-visual embedding",
     "Train a lecture-specific encoder that maps frame+caption pairs into a shared retrieval space with transcript text."),
    (ACCENT,"Stronger generators",
     "Evaluate GPT-4o and Claude Opus on the benchmark to measure the ceiling of RAG-constrained generation."),
    (RGBColor(0x88,0x44,0xAA), "Selective frame augmentation  (most actionable)",
     "Attach captions only to chunks where visual evidence is detected — avoids the –5.5% tIoU signal dilution on text-only multi-hop pairs."),
]
for i, (color, ftitle, fbody) in enumerate(fw):
    y = Inches(1.2) + i * Inches(1.2)
    label(s, Inches(0.35), y + Inches(0.06), Inches(0.52), Inches(0.45),
          str(i+1), size=Pt(15), fill=color)
    textbox(s, Inches(1.1), y + Inches(0.02), Inches(11.8), Inches(0.42),
            ftitle, size=Pt(17), bold=True, color=color)
    textbox(s, Inches(1.1), y + Inches(0.44), Inches(11.8), Inches(0.55),
            fbody, size=Pt(15), color=MID)
    if i < 4:
        divider(s, y + Inches(1.1), color=RGBColor(0xDD,0xDD,0xEE))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT)
title_bar(s, "Conclusion")

bullet_box(s, Inches(0.4), Inches(1.2), Inches(6.1), Inches(2.8),
    title="What we built",
    items=[
        "LectureBench — 810 temporally grounded QA pairs over 60 long lectures",
        "Only benchmark combining multi-hop + temporal spans + visual split for lecture-length video",
        "Open-source RAG pipeline with two configs and full evaluation suite",
    ], size=Pt(17), title_size=Pt(20))

bullet_box(s, Inches(0.4), Inches(4.1), Inches(6.1), Inches(2.6),
    title="What we found",
    items=[
        "Frame captions: +65% tIoU on visual-dependent questions",
        "Text-only questions: –3% (no degradation)",
        "Multi-hop retrieval is k-sensitive — larger k budgets help",
        "Visual blindness is a real, measurable failure mode in transcript-only RAG",
    ], size=Pt(17), title_size=Pt(20))

rect(s, Inches(6.8), Inches(1.15), Inches(6.15), Inches(4.6),
     fill=DARK, line=None)
textbox(s, Inches(7.0), Inches(1.5), Inches(5.75), Inches(0.55),
        "The one-slide answer", size=Pt(18), bold=True, color=C2)
textbox(s, Inches(7.0), Inches(2.1), Inches(5.75), Inches(3.0),
        "Adding Qwen2-VL-7B frame captions to lecture chunks improves temporally "
        "grounded retrieval by +65% on visual-dependent questions and leaves text-only "
        "questions unchanged — directly confirming that transcript-only RAG has a "
        "structural visual-evidence blind spot.",
        size=Pt(17), color=WHITE)
textbox(s, Inches(7.0), Inches(5.1), Inches(5.75), Inches(0.42),
        "github.com/NassimF/VQA-Benchmark  |  CC BY-NC-SA 4.0",
        size=Pt(13), color=C2, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, Inches(3.1), W, Inches(0.08), fill=C2)

textbox(s, Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.2),
        "Thank You", size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.5),
        "Questions?", size=Pt(28), color=C2, align=PP_ALIGN.CENTER)
textbox(s, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.45),
        "seyedehnasim.faridnia@utsa.edu",
        size=Pt(18), color=RGBColor(0xAA,0xBB,0xDD), align=PP_ALIGN.CENTER)
textbox(s, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.45),
        "github.com/NassimF/VQA-Benchmark",
        size=Pt(18), color=RGBColor(0xAA,0xBB,0xDD), align=PP_ALIGN.CENTER)

# stats bar
rect(s, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.2),
     fill=RGBColor(0x28,0x28,0x48), line=RGBColor(0x55,0x66,0x88))
for i, (num, lbl_) in enumerate([
    ("60", "lectures"), ("62.6h", "total video"), ("810", "QA pairs"), ("+65%", "visual tIoU")
]):
    xl = Inches(2.0) + i * Inches(2.55)
    textbox(s, xl, Inches(5.28), Inches(2.3), Inches(0.5),
            num, size=Pt(24), bold=True, color=C2, align=PP_ALIGN.CENTER)
    textbox(s, xl, Inches(5.78), Inches(2.3), Inches(0.35),
            lbl_, size=Pt(14), color=RGBColor(0xAA,0xBB,0xDD), align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────────────────────────
out = "/workspace/storage_nassim/VQA_Benchmark/slides/lecturebench_slides.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
